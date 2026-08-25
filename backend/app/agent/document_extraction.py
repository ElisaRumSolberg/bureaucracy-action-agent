import io
from dataclasses import dataclass

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
from pypdf.errors import PdfReadError

IMAGE_CONTENT_TYPES = {"image/jpeg", "image/png", "image/webp"}

TEXT_CONTENT_TYPES = {
    "application/pdf",
    "text/plain",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

SUPPORTED_CONTENT_TYPES = IMAGE_CONTENT_TYPES | TEXT_CONTENT_TYPES


class DocumentReadError(Exception):
    """Raised when the uploaded file's content couldn't be parsed."""


@dataclass
class ExtractedContent:
    text: str | None = None
    image_bytes: bytes | None = None
    image_mime_type: str | None = None

    @property
    def is_image(self) -> bool:
        return self.image_bytes is not None


def _extract_pdf_text(file_bytes: bytes) -> str:
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except PdfReadError as exc:
        raise DocumentReadError("We could not read this document.") from exc


def _extract_docx_text(file_bytes: bytes) -> str:
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as exc:  # noqa: BLE001
        raise DocumentReadError("We could not read this document.") from exc


def _extract_pptx_text(file_bytes: bytes) -> str:
    try:
        presentation = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)
    except Exception as exc:  # noqa: BLE001
        raise DocumentReadError("We could not read this document.") from exc


def _extract_txt_text(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")


def extract_document_content(file_bytes: bytes, content_type: str) -> ExtractedContent:
    """Reads an uploaded file into either plain text or (for images) raw
    bytes to hand to Gemini's vision input — no OCR library needed."""
    if content_type in IMAGE_CONTENT_TYPES:
        return ExtractedContent(image_bytes=file_bytes, image_mime_type=content_type)

    if content_type == "application/pdf":
        text = _extract_pdf_text(file_bytes)
    elif content_type == "text/plain":
        text = _extract_txt_text(file_bytes)
    elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = _extract_docx_text(file_bytes)
    elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = _extract_pptx_text(file_bytes)
    else:
        raise DocumentReadError(f"Unsupported file type: {content_type}")

    return ExtractedContent(text=text)
